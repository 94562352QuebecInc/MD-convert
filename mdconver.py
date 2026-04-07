#!/usr/bin/env python3
"""
mdconver.py — Batch document-to-Markdown converter using Multi-Provider APIs.
"""

import argparse
import base64
import datetime
import io
import json
import os
from typing import Dict, List, Optional, Tuple
import sys
import threading
import time
import traceback
import concurrent.futures

# ── Optional libraries ─────────────────────────────────────────────────────────
try:
    from pypdf import PdfReader, PdfWriter
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

import anthropic
import openai
import google.generativeai as genai

# ── Configuration ─────────────────────────────────────────────────────────────
BASE_DIR    = os.path.dirname(os.path.abspath(__file__))

def _load_api_keys() -> dict:
    keys = {
        "anthropic": os.environ.get("ANTHROPIC_API_KEY", "") or os.environ.get("AI_MODEL_MODEL_ANTHROPIC_KEY", ""),
        "gemini": os.environ.get("GEMINI_API_KEY", "") or os.environ.get("AI_MODEL_MODEL_GEMINI_API_KEY", ""),
        "openai": os.environ.get("OPENAI_API_KEY", "") or os.environ.get("AI_MODEL_MODEL_1773056808410_API_KEY", "") # Fallbacks from common env
    }
    
    env_path = os.path.join(BASE_DIR, ".env.local")
    if os.path.exists(env_path):
        with open(env_path) as f:
            for line in f:
                line = line.strip()
                if line.startswith("ANTHROPIC_API_KEY=") or line.startswith("AI_MODEL_MODEL_ANTHROPIC_KEY="):
                    if not keys["anthropic"]: keys["anthropic"] = line.split("=", 1)[1].strip().strip('"').strip("'")
                elif line.startswith("GEMINI_API_KEY=") or line.startswith("AI_MODEL_MODEL_GEMINI_API_KEY="):
                    if not keys["gemini"]: keys["gemini"] = line.split("=", 1)[1].strip().strip('"').strip("'")
                elif line.startswith("OPENAI_API_KEY="):
                    if not keys["openai"]: keys["openai"] = line.split("=", 1)[1].strip().strip('"').strip("'")
    return keys

API_KEYS = _load_api_keys()

MODEL_CONFIGS = {
    "claude-opus-4-6": {"provider": "anthropic", "pages_per_chunk": 20},
    "claude-sonnet-4-6":     {"provider": "anthropic", "pages_per_chunk": 20},
    "claude-haiku-4-5-20251001":    {"provider": "anthropic", "pages_per_chunk": 20},
    "gemini-3-flash-preview":             {"provider": "gemini",    "pages_per_chunk": 20},
    "gemini-3.1-flash-lite-preview":           {"provider": "gemini",    "pages_per_chunk": 20},
    "gpt-5.4-mini":               {"provider": "openai",    "pages_per_chunk": 20},
    "gpt-5.4-nano":               {"provider": "openai",    "pages_per_chunk": 20},
}

DEFAULT_MODEL = "claude-opus-4-6"

PDF_DIR     = os.path.join(BASE_DIR, "PDFs")
PROMPT_PATH = os.path.join(BASE_DIR, "Prompt", "prompt_file.md")
OUTPUT_DIR  = os.path.join(BASE_DIR, "Outputs")
LOG_DIR     = os.path.join(OUTPUT_DIR, "logs")

MAX_TOKENS_PER_PASS = 8000    # max output tokens commonly supported
MAX_PASSES          = 12      # safety cap on output continuation loops
MAX_CONNECTION_RETRIES = 3    # retries for transient connection errors
MAX_PDF_SIZE_MB     = 10      # PDFs larger than this are auto-split into chunks
CHUNK_PAGES         = 20      # pages per chunk when auto-splitting

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".xlsx", ".xls"}

# ── Logging ───────────────────────────────────────────────────────────────────

def _make_log_path(filename: str) -> str:
    ts        = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    stem      = os.path.splitext(filename)[0]
    safe_stem = "".join(c if c.isalnum() or c in " _-()." else "_" for c in stem).strip()
    return os.path.join(LOG_DIR, f"{ts}_{safe_stem}.log")


class RunLogger:
    def __init__(self, source_filename: str, model: str):
        self.source_filename = source_filename
        self.log_path        = _make_log_path(source_filename)
        self.start_time      = datetime.datetime.now()
        self.entries: List[Dict] = []
        self.total_input_tokens  = 0
        self.total_output_tokens = 0
        self._lines: List[str]   = []
        self.mode = "whole-pdf"
        self.lock = threading.RLock()
        self.model = model

    def log(self, msg: str):
        ts   = datetime.datetime.now().strftime("%H:%M:%S")
        line = f"[{ts}] {msg}"
        with self.lock:
            print(line, flush=True)
            self._lines.append(line)

    def record_call(self, label: str, pass_num: int,
                    input_tok: int, output_tok: int,
                    stop_reason: str, chars_generated: int):
        with self.lock:
            self.total_input_tokens  += input_tok
            self.total_output_tokens += output_tok
            entry = {
                "label":           label,
                "pass":            pass_num,
                "input_tokens":    input_tok,
                "output_tokens":   output_tok,
                "stop_reason":     stop_reason,
                "chars_generated": chars_generated,
            }
            self.entries.append(entry)
            self.log(
                f"  [{label} / pass {pass_num}] "
                f"in:{input_tok:,} tok  out:{output_tok:,} tok  "
                f"chars:{chars_generated:,}  stop:{stop_reason}"
            )

    def write(self, output_path: Optional[str], error: Optional[str] = None):
        os.makedirs(LOG_DIR, exist_ok=True)
        end_time = datetime.datetime.now()
        elapsed  = (end_time - self.start_time).total_seconds()

        with open(self.log_path, "w", encoding="utf-8") as f:
            f.write("=" * 72 + "\n")
            f.write(" mdconver.py — Run Log\n")
            f.write("=" * 72 + "\n\n")
            f.write(f"Source file   : {self.source_filename}\n")
            f.write(f"Model         : {self.model}\n")
            f.write(f"Mode          : {self.mode}\n")
            f.write(f"Run started   : {self.start_time.strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"Run finished  : {end_time.strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"Elapsed       : {elapsed:.1f}s\n")
            f.write(f"Output file   : {output_path or 'N/A (error)'}\n\n")

            f.write("-" * 72 + "\n")
            f.write(" Per-Call Token Summary\n")
            f.write("-" * 72 + "\n")
            f.write(f"{'Label':<30} {'Pass':>5}  {'Input':>10}  {'Output':>10}  {'Chars':>9}  Stop\n")
            f.write(f"{'-'*30} {'----':>5}  {'--------':>10}  {'--------':>10}  {'-----':>9}  ----\n")
            for e in self.entries:
                f.write(
                    f"{e['label']:<30} {e['pass']:>5}  {e['input_tokens']:>10,}  "
                    f"{e['output_tokens']:>10,}  {e['chars_generated']:>9,}  {e['stop_reason']}\n"
                )
            f.write(f"\n{'TOTAL':<30} {'':>5}  {self.total_input_tokens:>10,}  {self.total_output_tokens:>10,}\n\n")

            f.write("-" * 72 + "\n")
            f.write(" Console Output\n")
            f.write("-" * 72 + "\n")
            for line in self._lines:
                f.write(line + "\n")

            if error:
                f.write("\n" + "=" * 72 + "\n ERROR\n" + "=" * 72 + "\n")
                f.write(error + "\n")

        self.log(f"  Log saved → {self.log_path}")
        self.log(
            f"  TOTAL TOKENS — "
            f"input:{self.total_input_tokens:,}  "
            f"output:{self.total_output_tokens:,}  "
            f"combined:{self.total_input_tokens + self.total_output_tokens:,}"
        )

import re

# ── Multi-Extractor Baseline Framework ────────────────────────────────────────

try:
    from markitdown import MarkItDown
    HAS_MARKITDOWN = True
except ImportError:
    HAS_MARKITDOWN = False

try:
    import pymupdf4llm
    HAS_PYMUPDF4LLM = True
except ImportError:
    HAS_PYMUPDF4LLM = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False


def _normalize(text: str) -> str:
    """Normalize text for fuzzy comparison: lowercase, collapse whitespace, strip non-alphanum."""
    text = text.lower()
    text = re.sub(r'[^a-z0-9\s]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _to_meaningful_lines(raw: str) -> List[str]:
    """Split raw text into normalized meaningful lines (>=15 chars)."""
    lines = []
    for line in raw.split("\n"):
        normalized = _normalize(line)
        if len(normalized) >= 15:
            lines.append(normalized)
    return lines


# ── Individual Extractors ────────────────────────────────────────────────────

def _extract_markitdown(file_path: str) -> str:
    """Extract text using Microsoft MarkItDown."""
    if not HAS_MARKITDOWN:
        return ""
    try:
        md = MarkItDown()
        result = md.convert(file_path)
        return result.text_content or ""
    except Exception:
        return ""


def _extract_pymupdf4llm(file_path: str) -> str:
    """Extract text using pymupdf4llm (PyMuPDF-based, LLM-optimized markdown)."""
    if not HAS_PYMUPDF4LLM:
        return ""
    try:
        md_text = pymupdf4llm.to_markdown(
            file_path,
            page_chunks=False,
            write_images=False,
            show_progress=False,
        )
        return md_text or ""
    except Exception:
        return ""


def _extract_pdfplumber(file_path: str) -> str:
    """Extract text using pdfplumber (good for tables and structured layouts)."""
    if not HAS_PDFPLUMBER:
        return ""
    try:
        pages_text = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text(x_tolerance=2, y_tolerance=3) or ""
                # Also extract tables separately
                tables = page.extract_tables()
                table_md = ""
                for t_idx, table in enumerate(tables):
                    if not table:
                        continue
                    table_md += f"\n\n<!-- TABLE from pdfplumber page {i} -->\n"
                    for r_idx, row in enumerate(table):
                        cells = [str(c or "").strip() for c in row]
                        table_md += "| " + " | ".join(cells) + " |\n"
                        if r_idx == 0:
                            table_md += "|" + "|" .join(["---"] * len(cells)) + "|\n"
                pages_text.append(f"\n\n--- Page {i} ---\n\n{page_text}{table_md}")
        return "\n".join(pages_text)
    except Exception:
        return ""


# ── Multi-Extractor Orchestration ────────────────────────────────────────────

def _run_multi_extractor(file_path: str, logger: RunLogger) -> Tuple[str, Dict[str, str]]:
    """
    Run all available extractors, compare results, pick the best baseline.
    Returns (best_raw_text, dict_of_all_raw_texts_by_extractor_name).
    """
    ext = os.path.splitext(file_path)[1].lower()
    results: Dict[str, str] = {}

    if ext == ".pdf":
        logger.log("Running multi-extractor baseline (MarkItDown + pymupdf4llm + pdfplumber) …")

        # Run extractors in parallel
        with concurrent.futures.ThreadPoolExecutor(max_workers=3) as executor:
            futures = {}
            if HAS_MARKITDOWN:
                futures["markitdown"] = executor.submit(_extract_markitdown, file_path)
            if HAS_PYMUPDF4LLM:
                futures["pymupdf4llm"] = executor.submit(_extract_pymupdf4llm, file_path)
            if HAS_PDFPLUMBER:
                futures["pdfplumber"] = executor.submit(_extract_pdfplumber, file_path)

            for name, future in futures.items():
                try:
                    results[name] = future.result(timeout=120)
                except Exception as e:
                    logger.log(f"  ⚠️ {name} failed: {e}")
                    results[name] = ""
    else:
        # For non-PDF, only markitdown works
        if HAS_MARKITDOWN:
            results["markitdown"] = _extract_markitdown(file_path)

    # Score each extractor
    scores: Dict[str, dict] = {}
    for name, raw in results.items():
        lines = _to_meaningful_lines(raw)
        char_count = len(raw)
        line_count = len(lines)
        # Unique content — deduplicated meaningful lines
        unique_lines = len(set(lines))
        scores[name] = {
            "chars": char_count,
            "lines": line_count,
            "unique_lines": unique_lines,
            # Composite score: weight unique lines highest, then total chars
            "score": unique_lines * 100 + char_count,
        }
        logger.log(
            f"  📊 {name}: {char_count:,} chars, {line_count:,} meaningful lines, "
            f"{unique_lines:,} unique lines"
        )

    # Pick the best
    if not scores:
        logger.log("  ⚠️ No extractors produced output.")
        return "", results

    best_name = max(scores, key=lambda n: scores[n]["score"])
    logger.log(f"  ✅ Best extractor: {best_name} (score: {scores[best_name]['score']:,})")

    return results.get(best_name, ""), results


# ── Baseline Comparison & Coverage ───────────────────────────────────────────

def _run_baseline_comparison(
    file_path: str,
    md_text: str,
    logger: RunLogger,
    all_baselines: Optional[Dict[str, str]] = None,
) -> Tuple[str, str]:
    """
    Compare the best baseline text against the final output.
    Returns a coverage report string to append, and the best baseline raw text.
    """
    logger.log("Running completeness verification …")

    # Use the best baseline from multi-extractor
    if all_baselines:
        best_name = max(
            all_baselines,
            key=lambda n: len(_to_meaningful_lines(all_baselines[n])),
        )
        baseline_raw = all_baselines[best_name]
    else:
        baseline_raw = _extract_markitdown(file_path)
        best_name = "markitdown"

    baseline_lines = _to_meaningful_lines(baseline_raw)

    if not baseline_lines:
        logger.log("  ⚠️ Could not extract baseline text — skipping comparison.")
        return "", ""

    md_normalized = _normalize(md_text)

    found = 0
    missing = []
    for line in baseline_lines:
        if line in md_normalized:
            found += 1
        else:
            words = line.split()
            if len(words) >= 4:
                matches = sum(1 for w in words if w in md_normalized)
                if matches / len(words) >= 0.8:
                    found += 1
                    continue
            missing.append(line[:120])

    total = len(baseline_lines)
    coverage = (found / total * 100) if total > 0 else 100.0
    missing_ct = len(missing)

    logger.log(
        f"  Baseline ({best_name}): {total:,} lines  |  "
        f"Found: {found:,}  |  Missing: {missing_ct:,}  |  Coverage: {coverage:.1f}%"
    )

    # Also cross-check other extractors
    if all_baselines and len(all_baselines) > 1:
        for name, raw in all_baselines.items():
            if name == best_name:
                continue
            other_lines = _to_meaningful_lines(raw)
            if not other_lines:
                continue
            other_found = sum(1 for line in other_lines if line in md_normalized)
            other_cov = (other_found / len(other_lines) * 100) if other_lines else 100.0
            logger.log(f"  Cross-check ({name}): {len(other_lines):,} lines → {other_cov:.1f}% coverage")

    # Build the report
    report = []
    report.append("\n\n---\n")
    report.append("# APPENDIX — COMPLETENESS VERIFICATION")
    report.append("<!-- COMPLETENESS-CHECK -->")
    report.append("")
    report.append("> This section is auto-generated by the extraction pipeline. It compares text")
    report.append("> extracted via multiple Python extractors against the final Markdown output.")
    report.append("")

    # Extractor comparison table
    report.append("## Extractor Comparison")
    report.append("")
    report.append("| Extractor | Characters | Meaningful Lines | Unique Lines | Status |")
    report.append("|---|---|---|---|---|")
    for name, raw in (all_baselines or {}).items():
        lines = _to_meaningful_lines(raw)
        unique = len(set(lines))
        status = "✅ Best" if name == best_name else "⬜ Alt"
        report.append(f"| {name} | {len(raw):,} | {len(lines):,} | {unique:,} | {status} |")
    report.append("")

    # Coverage table
    report.append("## Coverage Results")
    report.append("")
    report.append("| Metric | Value |")
    report.append("|---|---|")
    report.append(f"| Reference extractor | {best_name} |")
    report.append(f"| Baseline text segments | {total:,} |")
    report.append(f"| Found in output | {found:,} |")
    report.append(f"| Potentially missing | {missing_ct:,} |")
    report.append(f"| **Coverage** | **{coverage:.1f}%** |")
    report.append("")

    if coverage >= 95:
        report.append("> ✅ **PASS** — Output covers ≥95% of the source text.")
    elif coverage >= 80:
        report.append(f"> ⚠️ **WARNING** — Output covers {coverage:.1f}%. Some content may be missing.")
    else:
        report.append(f"> ❌ **FAIL** — Output covers only {coverage:.1f}%. Significant content may be omitted.")

    if missing_ct > 0:
        report.append("")
        report.append("<details>")
        report.append(f"<summary>Show {min(missing_ct, 25)} potentially missing segments</summary>")
        report.append("")
        for i, line in enumerate(missing[:25]):
            report.append(f"{i+1}. `{line}`")
        if missing_ct > 25:
            report.append(f"")
            report.append(f"*... and {missing_ct - 25} more segments not shown.*")
        report.append("")
        report.append("</details>")

    return "\n".join(report), baseline_raw


# ── Baseline Page Splitter ───────────────────────────────────────────────────

def _split_baseline_by_pages(file_path: str, logger: RunLogger) -> Dict[int, str]:
    """
    Extract text page-by-page using the best available extractor.
    Returns dict mapping page_number (1-indexed) -> page_text.
    """
    pages: Dict[int, str] = {}
    ext = os.path.splitext(file_path)[1].lower()
    if ext != ".pdf":
        return pages

    # Use pymupdf4llm for page-level extraction (best quality)
    if HAS_PYMUPDF4LLM:
        try:
            chunks = pymupdf4llm.to_markdown(
                file_path,
                page_chunks=True,
                write_images=False,
                show_progress=False,
            )
            for chunk in chunks:
                pg = chunk.get("metadata", {}).get("page_number", 1)
                text = chunk.get("text", "")
                if text.strip():  # Skip empty pages
                    pages[pg] = text
            if pages:
                logger.log(f"  Page-split baseline: {len(pages)} pages via pymupdf4llm")
                return pages
        except Exception:
            pass

    # Fallback to pdfplumber
    if HAS_PDFPLUMBER:
        try:
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text(x_tolerance=2, y_tolerance=3) or ""
                    pages[i] = text
            if pages:
                logger.log(f"  Page-split baseline: {len(pages)} pages via pdfplumber")
                return pages
        except Exception:
            pass

    # Fallback to pypdf
    if HAS_PYPDF:
        try:
            reader = PdfReader(file_path)
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                pages[i] = text
            if pages:
                logger.log(f"  Page-split baseline: {len(pages)} pages via pypdf")
        except Exception:
            pass

    return pages


# ── File helpers ──────────────────────────────────────────────────────────────

def load_prompt() -> str:
    with open(PROMPT_PATH, "r", encoding="utf-8") as f:
        return f.read()

def load_pdf_b64(path: str) -> str:
    with open(path, "rb") as f:
        return base64.standard_b64encode(f.read()).decode("utf-8")

def get_pdf_page_count(path: str) -> int:
    if not HAS_PYPDF:
        raise RuntimeError("pypdf not installed. Run: pip install pypdf")
    return len(PdfReader(path).pages)

def extract_pages_b64(path: str, start: int, end: int) -> str:
    reader = PdfReader(path)
    writer = PdfWriter()
    for i in range(start, min(end, len(reader.pages))):
        writer.add_page(reader.pages[i])
    buf = io.BytesIO()
    writer.write(buf)
    return base64.standard_b64encode(buf.getvalue()).decode("utf-8")

def extract_pages_to_file(path: str, start: int, end: int, output_path: str) -> str:
    """Write a subset of PDF pages to a new file on disk. Returns output_path."""
    reader = PdfReader(path)
    writer = PdfWriter()
    for i in range(start, min(end, len(reader.pages))):
        writer.add_page(reader.pages[i])
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path

def _pdf_to_text(path: str, start: int=0, end: int=None) -> str:
    if not HAS_PYPDF:
        raise RuntimeError("pypdf not installed.")
    reader = PdfReader(path)
    lines = []
    end = end if end else len(reader.pages)
    for i in range(start, min(end, len(reader.pages))):
        text = reader.pages[i].extract_text()
        if text: lines.append(text)
    return "\n".join(lines)

def pptx_to_text(path: str) -> str:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed.")
    prs   = Presentation(path)
    lines = []
    for n, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- SLIDE {n} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
            if shape.shape_type == 19:
                for row in shape.table.rows:
                    lines.append("  |  ".join(c.text.strip() for c in row.cells))
    return "\n".join(lines)

def docx_to_text(path: str) -> str:
    if not HAS_DOCX:
        raise RuntimeError("python-docx not installed.")
    doc = docx.Document(path)
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in doc.tables.rows:
            lines.append("  |  ".join(c.text.strip() for c in row.cells))
    return "\n".join(lines)

def xlsx_to_text(path: str) -> str:
    if not HAS_XLSX:
        raise RuntimeError("openpyxl not installed.")
    wb    = openpyxl.load_workbook(path, data_only=True)
    lines = []
    for name in wb.sheetnames:
        lines.append(f"\n=== SHEET: {name} ===\n")
        for row in wb[name].iter_rows(values_only=True):
            row_text = "  |  ".join("" if c is None else str(c) for c in row)
            if row_text.strip("|").strip():
                lines.append(row_text)
    return "\n".join(lines)

def _build_doc_block(file_path: str, provider: str) -> list:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        if provider == "anthropic":
            return [{"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": load_pdf_b64(file_path)}}]
        elif provider == "gemini":
            import google.generativeai as genai
            return [{"mime_type": "application/pdf", "data": load_pdf_b64(file_path)}]
        elif provider == "openai":
            # OpenAI does not natively support base64 PDF in standard chat completions, extracting text
            return [{"type": "text", "text": f"[Source: PDF Document]\n\n{_pdf_to_text(file_path)}"}]
    
    # Generic format handling for PPTX, DOCX, XLSX
    if ext == ".pptx":   text, label = pptx_to_text(file_path), "PowerPoint"
    elif ext == ".docx": text, label = docx_to_text(file_path),  "Word document"
    elif ext in (".xlsx", ".xls"): text, label = xlsx_to_text(file_path), "Excel spreadsheet"
    else: raise ValueError(f"Unsupported type: {ext}")
    return [{"type": "text", "text": f"[Source: {label}]\n\n{text}"}]

def _build_chunk_doc_block(file_path: str, start: int, end: int, provider: str) -> list:
    if provider == "anthropic":
        return [{"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": extract_pages_b64(file_path, start, end)}}]
    elif provider == "gemini":
        return [{"mime_type": "application/pdf", "data": extract_pages_b64(file_path, start, end)}]
    elif provider == "openai":
        return [{"type": "text", "text": f"[Source: PDF Document Chunk pp.{start+1}-{end}]\n\n{_pdf_to_text(file_path, start, end)}"}]
    return []

# ── Prompts ───────────────────────────────────────────────────────────────────

CONTINUATION_INSTRUCTION = """\
Your output was cut off because you reached the output token limit.
CONTINUE from EXACTLY where you left off. Rules:
- Do NOT repeat any content already written above.
- Do NOT restart from the beginning. Do NOT add any preamble.
- Jump straight back into the content and continue until complete.
"""

def _chunk_enhance_prompt(start_page: int, end_page: int, total_pages: int, baseline_text: str) -> str:
    return f"""\
You are enhancing pages {start_page}–{end_page} of a {total_pages}-page document.

BELOW is the raw Python-extracted text for these pages. This text is COMPLETE but lacks
Rich formatting. You have the original PDF pages attached for visual reference.

Your task is to ENHANCE this baseline text — NOT re-extract it. Rules:

1. **PRESERVE ALL TEXT**: You must keep every word from the baseline. Do NOT omit, summarize,
   or paraphrase any content. The baseline text is the ground truth.

2. **ADD STRUCTURAL TAGS**:
   - Insert `<!-- PAGE: N -->` before each page and `---` between pages.
   - Wrap every table: `<!-- TABLE-START: <title> -->` / `<!-- TABLE-END: <title> -->`
   - Add table heading: `#### [TABLE: <title>]` with `**Page:** N`
   - Wrap every chart: `<!-- CHART-START: <title> -->` / `<!-- CHART-END: <title> -->`
   - Wrap images: `<!-- IMAGE: <description> -->`
   - Add navigation tags: `<!-- FUND-LEVEL -->`, `<!-- PORTFOLIO-COMPANY: <name> -->`,
     `<!-- FINANCIALS -->`, `<!-- KPI: <metric> -->`, `<!-- MANAGER-LETTER -->`,
     `<!-- EXIT: <name> -->`, `<!-- NEW-INVESTMENT: <name> -->`, `<!-- WARNING: <reason> -->`

3. **FORMAT TABLES**: Convert raw tabular text into proper Markdown tables.
   Add 2–4 sentence explanations below each table.

4. **DESCRIBE CHARTS/IMAGES**: Use the PDF visual to describe charts the baseline cannot capture.
   Include chart type, axes, key data points, and 2–3 sentence context.

5. **TAG REPETITIVE CONTENT**: Identify and tag repeated elements:
   - `<!-- REPEATED-HEADER: <text> -->` for recurring page headers
   - `<!-- REPEATED-FOOTER: <text> -->` for recurring page footers/disclaimers  
   - `<!-- PAGE-NUMBER: N -->` for standalone page numbers

6. **MERGE not REPLACE**: If the baseline has garbled text (OCR artifacts, broken lines),
   clean it up. But if the baseline has valid text that doesn't appear in the PDF visual
   (e.g., metadata), keep it anyway.

Begin your output directly with `<!-- PAGE: {start_page} -->` and continue through page {end_page}.

--- BASELINE TEXT FOR PAGES {start_page}-{end_page} ---

{baseline_text}

--- END BASELINE TEXT ---
"""


def _synthesis_prompt(prompt_text: str) -> str:
    return f"""\
{prompt_text}

────────────────────────────────────────────────────────────
IMPORTANT INSTRUCTION FOR THIS CALL:

The SECTION 2 full-document Markdown baseline has already been provided to you.
It is complete — do NOT re-extract or re-write it.

Your task is ONLY to produce the following and nothing else:
1. The document header block
2. SECTION 0 — DOCUMENT NAVIGATOR
3. SECTION 1 — EXTRACTED SCHEMA DATA
4. SECTION 3 — STRUCTURAL TAGS (JSON block)

Output these sections exactly as requested in the prompt.
Do NOT reproduce Section 2.
────────────────────────────────────────────────────────────
"""

# ── Core API call helper ──────────────────────────────────────────────────────

def _stream_with_continuation(
    messages: list,
    label: str,
    logger: RunLogger,
    model: str,
    provider: str
) -> str:
    """Multi-provider streaming logic."""
    accumulated = ""
    pass_num    = 0

    if provider == "anthropic":
        client = anthropic.Anthropic(api_key=API_KEYS.get("anthropic"))
    elif provider == "openai":
        client = openai.OpenAI(api_key=API_KEYS.get("openai"))
    elif provider == "gemini":
        genai.configure(api_key=API_KEYS.get("gemini"))
        client = genai.GenerativeModel(model)

    while pass_num < MAX_PASSES:
        pass_num += 1
        logger.log(f"  [{label}] Starting pass {pass_num} …")

        _stop_hb = threading.Event()
        def _heartbeat(lbl=label, pn=pass_num):
            count = 0
            while not _stop_hb.wait(10):
                count += 1
                print(f"  [{lbl} / pass {pn}] … API call in progress ({count * 10}s elapsed) …", flush=True)
        hb = threading.Thread(target=_heartbeat, daemon=True)
        hb.start()

        pass_text     = ""
        stop_reason   = None
        input_tokens  = 0
        output_tokens = 0

        try:
            if provider == "anthropic":
                with client.messages.stream(
                    model=model, max_tokens=MAX_TOKENS_PER_PASS, messages=messages,
                ) as stream:
                    chunk_count = 0
                    for text_chunk in stream.text_stream:
                        pass_text += text_chunk
                        chunk_count += 1
                        if chunk_count % 200 == 0: logger.log(f"  [{label}] … streaming: {len(accumulated) + len(pass_text):,} chars …")
                    final_msg     = stream.get_final_message()
                    stop_reason   = final_msg.stop_reason
                    input_tokens  = getattr(final_msg.usage, 'input_tokens', 0)
                    output_tokens = getattr(final_msg.usage, 'output_tokens', 0)
            elif provider == "openai":
                # Convert messages format if needed
                oai_msgs = []
                for m in messages:
                    c = m.get("content", "")
                    if isinstance(c, list):
                        txt = ""
                        for b in c:
                            if b.get("type") == "text": txt += b.get("text", "") + "\n"
                        oai_msgs.append({"role": m["role"], "content": txt})
                    else:
                        oai_msgs.append({"role": m["role"], "content": c})
                
                response = client.chat.completions.create(
                    model=model, messages=oai_msgs, stream=True
                )
                for chunk in response:
                    if chunk.choices and chunk.choices[0].delta.content:
                        pass_text += chunk.choices[0].delta.content
                # we don't get exact token counts from standard streaming cleanly, we'll estimate token/char ratio (roughly 1 tok = 4 chars)
                input_tokens = sum(len(str(m["content"])) for m in oai_msgs) // 4
                output_tokens = len(pass_text) // 4
                
                # Retrieve final string stop reason by issuing a secondary small check or assuming stream completion
                stop_reason = chunk.choices[0].finish_reason if chunk.choices else "stop"
                if stop_reason == "length": stop_reason = "max_tokens"
                if stop_reason in ["stop", "None", None]: stop_reason = "end_turn"
            elif provider == "gemini":
                # Gemini expects alternating user/model
                gemini_msgs = []
                system_instruction = ""
                for m in messages:
                    role = "user" if m["role"] == "user" else "model"
                    content_parts = []
                    
                    if isinstance(m["content"], list):
                        for part in m["content"]:
                            if isinstance(part, str):
                                # Plain string part — add directly as text
                                content_parts.append(part)
                            elif isinstance(part, dict):
                                if "mime_type" in part:
                                    content_parts.append({"mime_type": part["mime_type"], "data": part["data"]})
                                elif part.get("type") == "text":
                                    content_parts.append(part.get("text", ""))
                                else:
                                    content_parts.append(str(part))
                            else:
                                content_parts.append(str(part))
                    else:
                        content_parts = [m["content"]]
                        
                    gemini_msgs.append({"role": role, "parts": content_parts})

                response = client.generate_content(gemini_msgs, stream=True)
                for chunk in response:
                    if chunk.text: pass_text += chunk.text
                usage = response.usage_metadata
                input_tokens = usage.prompt_token_count if usage else 0
                output_tokens = usage.candidates_token_count if usage else 0
                # Map finish reasons
                candidate = response.candidates[0] if response.candidates else None
                if candidate and candidate.finish_reason == 2: # MAX_TOKENS
                    stop_reason = "max_tokens"
                else: stop_reason = "end_turn"

        except Exception as e:
            _stop_hb.set()
            err_str = str(e).lower()
            if "content filtering" in err_str or "blocked" in err_str or "safety" in err_str:
                logger.log(f"  [{label}] ⚠️ Content filter triggered: {e}")
                pass_text = f"\n\n> [!CAUTION]\n> Content for this section was blocked by provider content filtering.\n\n"
                stop_reason = "content_filter"
            else:
                if not hasattr(_stream_with_continuation, '_retry_count'):
                    _stream_with_continuation._retry_count = 0
                
                if "413" in err_str or "too_large" in err_str or "request_too_large" in err_str or "payload" in err_str:
                    if any(isinstance(c, dict) and (c.get("type") == "document" or "mime_type" in c) for m in messages if isinstance(m.get("content"), list) for c in m["content"]):
                        logger.log(f"  [{label}] ⚠️ 413 Payload Too Large. Stripping visual PDF to retry text-only...")
                        for m in messages:
                            if isinstance(m.get("content"), list):
                                m["content"] = [c for c in m["content"] if not (isinstance(c, dict) and (c.get("type") == "document" or "mime_type" in c))]
                        pass_num -= 1
                        time.sleep(1)
                        continue

                _stream_with_continuation._retry_count += 1
                if _stream_with_continuation._retry_count <= MAX_CONNECTION_RETRIES:
                    wait = 2 ** _stream_with_continuation._retry_count
                    logger.log(f"  [{label}] Connection error: {e}. Retrying in {wait}s… (attempt {_stream_with_continuation._retry_count}/{MAX_CONNECTION_RETRIES})")
                    time.sleep(wait)
                    pass_num -= 1
                    continue
                else:
                    _stream_with_continuation._retry_count = 0
                    raise e
        finally:
            _stop_hb.set()

        _stream_with_continuation._retry_count = 0

        logger.record_call(label, pass_num, input_tokens, output_tokens,
                           stop_reason, len(pass_text))
        accumulated += pass_text

        if stop_reason == "end_turn":
            logger.log(f"  [{label}] Complete after {pass_num} pass(es).")
            break
        elif stop_reason == "max_tokens":
            logger.log(f"  [{label}] Output truncated — running continuation pass …")
            messages = messages + [
                {"role": "assistant" if provider != "gemini" else "model", "content": accumulated},
                {"role": "user",      "content": CONTINUATION_INSTRUCTION},
            ]
        else:
            logger.log(f"  [{label}] Unexpected stop_reason: {stop_reason}. Stopping.")
            break

    return accumulated

# ── JSON Patch Application ────────────────────────────────────────────────────

def _apply_json_patches(llm_output: str, baseline_text: str, logger: RunLogger, label: str = "") -> str:
    """Parse JSON patch array from LLM output and inject tags into baseline text."""
    import re as _re
    patched = baseline_text
    json_match = _re.search(r'```json\s*(\[.*?\])\s*```', llm_output, _re.DOTALL)
    if json_match:
        try:
            patches = json.loads(json_match.group(1))
            success_count = 0
            for patch in patches:
                tag = patch.get("tag", "")
                anchor = patch.get("insert_before_exact_string", "")
                if tag and anchor and anchor in patched:
                    patched = patched.replace(anchor, f"{tag}\n{anchor}", 1)
                    success_count += 1
            logger.log(f"  {label}Applied {success_count}/{len(patches)} JSON patches.")
        except Exception as e:
            logger.log(f"  {label}⚠️ Failed to parse/apply JSON patches: {e}")
    else:
        logger.log(f"  {label}⚠️ No JSON patch block found in LLM output.")
    return patched

# ── Unified Synthesis Strategy ───────────────────────────────────────────────────

def _run_synthesis(
    file_path: str,
    prompt_text: str,
    logger: RunLogger,
    model: str,
    provider: str,
    baseline_text: str
) -> str:
    """Single-pass synthesis: sends full PDF + baseline to LLM, gets Sections 0/1/3, patches Section 2."""
    logger.log("Strategy: Unified Synthesis + JSON Patch (Zero-Data-Loss Pipeline)")
    
    doc_block = _build_doc_block(file_path, provider)
    synth_prompt = _synthesis_prompt(prompt_text)
    
    # Cap to avoid token bloat
    baseline_for_synth = baseline_text[:500_000]
    if len(baseline_text) > 500_000:
         logger.log("  Section 2 truncated to 500,000 chars for synthesis pass.")
         
    if provider == "gemini":
        messages = [{"role": "user", "content": doc_block + [{"type": "text", "text": f"{synth_prompt}\n\n--- BASELINE ---\n\n{baseline_for_synth}"}]}]
    else:
        messages = [{"role": "user", "content": doc_block + [{"type": "text", "text": f"{synth_prompt}\n\n--- BASELINE ---\n\n{baseline_for_synth}"}]}]

    output_013 = _stream_with_continuation(messages, "synthesis", logger, model, provider)

    sections_header = output_013.rstrip()
    patched_section2 = _apply_json_patches(output_013, baseline_text, logger)

    final_md = (
        sections_header
        + "\n\n---\n\n"
        + "# SECTION 2 — FULL DOCUMENT MARKDOWN\n\n"
        + "> This section contains the complete text of the source document.\n"
        + "> The text was extracted via Python (pymupdf4llm/pdfplumber/markitdown)\n"
        + "> and mechanically enhanced with structural tags.\n"
        + "> Use the inline `<!-- TAGS -->` to navigate to specific sections quickly.\n\n"
        + patched_section2
    )

    return final_md

# ── Chunked Synthesis Strategy ────────────────────────────────────────────────

def _run_chunked_synthesis(
    file_path: str,
    prompt_text: str,
    logger: RunLogger,
    model: str,
    provider: str,
    baseline_text: str,
    baseline_pages: Dict[int, str]
) -> str:
    """Split a large PDF into chunks, process each independently, then merge."""
    import tempfile
    total_pages = get_pdf_page_count(file_path)
    logger.log(f"Strategy: Chunked Synthesis (Zero-Data-Loss) — {total_pages} pages, {CHUNK_PAGES} per chunk")

    # ── Step 1: Split PDF into temp chunk files ──
    chunks = []   # list of (start_page_1idx, end_page_1idx, chunk_pdf_path, chunk_baseline)
    tmp_dir = tempfile.mkdtemp(prefix="mdconv_chunks_")
    
    for start_0 in range(0, total_pages, CHUNK_PAGES):
        end_0 = min(start_0 + CHUNK_PAGES, total_pages)
        start_1 = start_0 + 1
        end_1 = end_0
        
        chunk_path = os.path.join(tmp_dir, f"chunk_pp{start_1}-{end_1}.pdf")
        extract_pages_to_file(file_path, start_0, end_0, chunk_path)
        
        # Gather baseline text for this chunk's pages
        chunk_baseline_parts = []
        for pg in range(start_1, end_1 + 1):
            pg_text = baseline_pages.get(pg, "")
            if pg_text.strip():
                chunk_baseline_parts.append(f"\n--- Page {pg} ---\n{pg_text}")
        chunk_baseline = "\n".join(chunk_baseline_parts) if chunk_baseline_parts else "(No baseline text available)"
        
        chunks.append((start_1, end_1, chunk_path, chunk_baseline))
        chunk_size_mb = os.path.getsize(chunk_path) / (1024 * 1024)
        logger.log(f"  Created chunk pp{start_1}-{end_1}: {chunk_size_mb:.1f} MB, {len(chunk_baseline):,} chars baseline")

    # ── Step 2: Process each chunk independently ──
    all_patches: list = []  # collect all JSON patches across chunks
    sections_header = ""    # will use the first chunk's Section 0/1 as the master

    for idx, (start_1, end_1, chunk_path, chunk_baseline) in enumerate(chunks):
        label = f"chunk-{idx+1} (pp{start_1}-{end_1})"
        logger.log(f"Processing {label} …")
        
        doc_block = _build_doc_block(chunk_path, provider)
        synth_prompt = _synthesis_prompt(prompt_text)
        
        baseline_for_synth = chunk_baseline[:500_000]
        
        if provider == "gemini":
            messages = [{"role": "user", "content": doc_block + [{"type": "text", "text": f"{synth_prompt}\n\n--- BASELINE (Pages {start_1}-{end_1}) ---\n\n{baseline_for_synth}"}]}]
        else:
            messages = [{"role": "user", "content": doc_block + [{"type": "text", "text": f"{synth_prompt}\n\n--- BASELINE (Pages {start_1}-{end_1}) ---\n\n{baseline_for_synth}"}]}]

        chunk_output = _stream_with_continuation(messages, label, logger, model, provider)
        
        # Capture Section 0/1 from first chunk only (it sees the document start)
        if idx == 0:
            sections_header = chunk_output.rstrip()
        
        # Extract and collect JSON patches from this chunk
        import re as _re
        json_match = _re.search(r'```json\s*(\[.*?\])\s*```', chunk_output, _re.DOTALL)
        if json_match:
            try:
                chunk_patches = json.loads(json_match.group(1))
                all_patches.extend(chunk_patches)
                logger.log(f"  [{label}] Extracted {len(chunk_patches)} JSON patches.")
            except Exception as e:
                logger.log(f"  [{label}] ⚠️ Failed to parse JSON patches: {e}")
        else:
            logger.log(f"  [{label}] ⚠️ No JSON patch block found.")
        
        # Rate-limit between chunks
        if idx < len(chunks) - 1:
            logger.log("  Waiting 5s before next chunk …")
            time.sleep(5)

    # ── Step 3: Apply all patches to the FULL baseline ──
    patched_section2 = baseline_text
    success_count = 0
    for patch in all_patches:
        tag = patch.get("tag", "")
        anchor = patch.get("insert_before_exact_string", "")
        if tag and anchor and anchor in patched_section2:
            patched_section2 = patched_section2.replace(anchor, f"{tag}\n{anchor}", 1)
            success_count += 1
    logger.log(f"Applied {success_count}/{len(all_patches)} total JSON patches to full baseline.")

    # ── Step 4: Assemble final document ──
    final_md = (
        sections_header
        + "\n\n---\n\n"
        + "# SECTION 2 — FULL DOCUMENT MARKDOWN\n\n"
        + "> This section contains the complete text of the source document.\n"
        + "> The text was extracted via Python (pymupdf4llm/pdfplumber/markitdown)\n"
        + "> and mechanically enhanced with structural tags.\n"
        + "> Use the inline `<!-- TAGS -->` to navigate to specific sections quickly.\n\n"
        + patched_section2
    )

    # Clean up temp files
    import shutil
    try:
        shutil.rmtree(tmp_dir, ignore_errors=True)
    except Exception:
        pass

    return final_md

# ── Main entry per file ───────────────────────────────────────────────────────

def run_single_file(
    file_path: str,
    prompt_text: str,
    logger: RunLogger,
    model: str,
    best_baseline: str,
) -> str:
    config = MODEL_CONFIGS.get(model, {"provider": "anthropic", "pages_per_chunk": 20})
    provider = config.get("provider", "anthropic")
    
    ext = os.path.splitext(file_path)[1].lower()
    file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
    
    if ext == ".pdf" and file_size_mb > MAX_PDF_SIZE_MB:
        logger.log(f"📦 PDF size ({file_size_mb:.1f} MB) exceeds {MAX_PDF_SIZE_MB} MB limit — auto-splitting into {CHUNK_PAGES}-page chunks.")
        baseline_pages = _split_baseline_by_pages(file_path, logger)
        return _run_chunked_synthesis(file_path, prompt_text, logger, model, provider, best_baseline, baseline_pages)
    
    return _run_synthesis(file_path, prompt_text, logger, model, provider, best_baseline)

# ── Batch runner ──────────────────────────────────────────────────────────────

def convert_file(file_path: str, prompt_text: str, model: str) -> dict:
    filename    = os.path.basename(file_path)
    stem        = os.path.splitext(filename)[0]
    output_path = os.path.join(OUTPUT_DIR, stem + ".md")
    logger      = RunLogger(filename, model)

    logger.log("=" * 62)
    logger.log(f"Processing : {filename}")
    logger.log(f"Model      : {model}")
    logger.log(f"Output     : {output_path}")

    try:
        os.makedirs(OUTPUT_DIR, exist_ok=True)

        # ── Phase 1: Multi-extractor baseline ──
        best_baseline, all_baselines = _run_multi_extractor(file_path, logger)

        # ── Phase 2: LLM enhancement pipeline ──
        md_text = run_single_file(file_path, prompt_text, logger, model, best_baseline)

        # ── Phase 3: Completeness verification ──
        comparison_report, ref_baseline = _run_baseline_comparison(
            file_path, md_text, logger, all_baselines
        )
        if comparison_report:
            md_text += comparison_report

        # ── Save all baselines ──
        for ext_name, raw in all_baselines.items():
            if not raw:
                continue
            suffix = f"_{ext_name.upper()}.md"
            ext_path = os.path.join(OUTPUT_DIR, stem + suffix)
            with open(ext_path, "w", encoding="utf-8") as f:
                f.write(raw)
            logger.log(f"📄 Saved {ext_name} baseline → {ext_path}  ({len(raw):,} chars)")

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(md_text)

        logger.log(f"✅ Saved → {output_path}  ({len(md_text):,} chars)")
        logger.write(output_path)
        return {
            "file": filename, "output": output_path, "success": True,
            "mode": logger.mode, "passes": len(logger.entries),
            "total_input_tokens":  logger.total_input_tokens,
            "total_output_tokens": logger.total_output_tokens,
            "log": logger.log_path,
        }

    except Exception as exc:
        err = traceback.format_exc()
        logger.log(f"❌ ERROR: {exc}")
        logger.write(None, error=err)
        return {
            "file": filename, "success": False,
            "error": str(exc), "log": logger.log_path,
        }

def main():
    parser = argparse.ArgumentParser(description="Batch document-to-Markdown converter")
    parser.add_argument("--model", type=str, default=DEFAULT_MODEL, help="Model to use")
    parser.add_argument("files", nargs="*", help="Files to process (basenames in PDFs/ or absolute paths)")
    args = parser.parse_args()

    prompt_text = load_prompt()
    print(f"Prompt loaded: {len(prompt_text):,} chars", flush=True)

    if args.files:
        files = [f if os.path.isabs(f) else os.path.join(PDF_DIR, f) for f in args.files]
    else:
        files = sorted([
            os.path.join(PDF_DIR, f) for f in os.listdir(PDF_DIR)
            if os.path.splitext(f)[1].lower() in SUPPORTED_EXTENSIONS
        ])

    if not files:
        print("No supported files found. Nothing to do.", flush=True)
        return

    print(f"Files to process: {len(files)}", flush=True)
    for f in files:
        print(f"  • {os.path.basename(f)}", flush=True)
    print(flush=True)

    results             = []
    grand_input_tokens  = 0
    grand_output_tokens = 0

    for file_path in files:
        result = convert_file(file_path, prompt_text, args.model)
        results.append(result)
        grand_input_tokens  += result.get("total_input_tokens",  0)
        grand_output_tokens += result.get("total_output_tokens", 0)
        print(flush=True)

    print("=" * 62, flush=True)
    print("SESSION SUMMARY", flush=True)
    print("=" * 62, flush=True)
    for r in results:
        icon = "✅" if r["success"] else "❌"
        info = (
            f"mode:{r.get('mode','?')}  "
            f"in:{r.get('total_input_tokens', 0):,}  "
            f"out:{r.get('total_output_tokens', 0):,}  "
            f"calls:{r.get('passes', '?')}"
        ) if r["success"] else f"ERROR: {r.get('error', '')}"
        print(f"  {icon}  {r['file']:<55} {info}", flush=True)
        print(f"       Log → {r.get('log', 'N/A')}", flush=True)

    print(flush=True)
    print(f"  Grand total input tokens  : {grand_input_tokens:,}",  flush=True)
    print(f"  Grand total output tokens : {grand_output_tokens:,}", flush=True)
    print(f"  Grand total combined      : {grand_input_tokens + grand_output_tokens:,}", flush=True)
    print("=" * 62, flush=True)

    ts          = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    session_log = os.path.join(LOG_DIR, f"{ts}_session_summary.json")
    os.makedirs(LOG_DIR, exist_ok=True)
    with open(session_log, "w", encoding="utf-8") as f:
        json.dump({
            "timestamp":           ts,
            "model":               args.model,
            "files_processed":     len(results),
            "grand_input_tokens":  grand_input_tokens,
            "grand_output_tokens": grand_output_tokens,
            "grand_total_tokens":  grand_input_tokens + grand_output_tokens,
            "results":             results,
        }, f, indent=2)
    print(f"  Session JSON → {session_log}", flush=True)

if __name__ == "__main__":
    main()
